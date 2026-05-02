#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPT Subtitle Generator for Mumble Jumble (乱民全讲)
Generates PowerPoint subtitles from translation JSON files
Created by Claude Code

Usage:
    python3 generate_subtitles_ppt.py scene_1-4_translation.json
    python3 generate_subtitles_ppt.py translations.json --all
"""

import json
import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Configuration
MAX_CHARS_PER_SLIDE = 100  # Maximum characters per slide
SLIDE_WIDTH = Inches(10)  # 16:9 aspect ratio
SLIDE_HEIGHT = Inches(5.625)
BG_COLOR = RGBColor(0, 0, 0)  # Black background
TEXT_COLOR = RGBColor(255, 255, 255)  # White text
CHARACTER_COLOR = RGBColor(200, 200, 200)  # Light gray for character names


def create_presentation():
    """Create a new PowerPoint presentation with black background"""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    return prs


def add_subtitle_slide(prs, character_pinyin, translation, scene_id=None):
    """
    Add a subtitle slide to the presentation

    Args:
        prs: Presentation object
        character_pinyin: Character name in pinyin (e.g., "XING")
        translation: English translation text
        scene_id: Optional scene identifier (e.g., "1-4")
    """
    # Add blank slide
    blank_slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_slide_layout)

    # Set black background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

    # Add character name (top, smaller font)
    char_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5),
        Inches(9), Inches(0.8)
    )
    char_frame = char_box.text_frame
    char_frame.text = character_pinyin
    char_para = char_frame.paragraphs[0]
    char_para.font.size = Pt(24)
    char_para.font.bold = True
    char_para.font.color.rgb = CHARACTER_COLOR
    char_para.alignment = PP_ALIGN.LEFT

    # Add translation (center, large font)
    trans_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.5),
        Inches(9), Inches(3)
    )
    trans_frame = trans_box.text_frame
    trans_frame.text = translation
    trans_frame.word_wrap = True
    trans_para = trans_frame.paragraphs[0]
    trans_para.font.size = Pt(36)
    trans_para.font.color.rgb = TEXT_COLOR
    trans_para.alignment = PP_ALIGN.CENTER

    # Add scene info (bottom right corner, small)
    if scene_id:
        scene_box = slide.shapes.add_textbox(
            Inches(8), Inches(5),
            Inches(1.5), Inches(0.5)
        )
        scene_frame = scene_box.text_frame
        scene_frame.text = f"Scene {scene_id}"
        scene_para = scene_frame.paragraphs[0]
        scene_para.font.size = Pt(14)
        scene_para.font.color.rgb = RGBColor(100, 100, 100)
        scene_para.alignment = PP_ALIGN.RIGHT


def split_long_text(text, max_length=MAX_CHARS_PER_SLIDE):
    """
    Split long text into multiple parts at natural break points

    Breaks at sentence boundaries (. ! ?) when possible, otherwise at phrases

    Args:
        text: Text to split
        max_length: Maximum characters per part

    Returns:
        List of text segments
    """
    if len(text) <= max_length:
        return [text]

    # Try to split at sentence boundaries first
    sentence_endings = ['. ', '! ', '? ']
    parts = []
    remaining = text

    while len(remaining) > max_length:
        # Find the best split point within max_length
        best_split = -1

        # Look for sentence endings
        for ending in sentence_endings:
            pos = remaining.rfind(ending, 0, max_length)
            if pos > best_split:
                best_split = pos + len(ending)

        # If no sentence ending found, try commas or dashes
        if best_split == -1:
            for punct in [', ', '—', '; ', ': ']:
                pos = remaining.rfind(punct, 0, max_length)
                if pos > best_split:
                    best_split = pos + len(punct)

        # If still no good split point, split at last space
        if best_split == -1:
            best_split = remaining.rfind(' ', 0, max_length)
            if best_split == -1:
                best_split = max_length

        # Extract part and continue
        parts.append(remaining[:best_split].strip())
        remaining = remaining[best_split:].strip()

    # Add remaining text
    if remaining:
        parts.append(remaining)

    return parts


def should_group_lines(current_line, next_line, current_length):
    """
    Determine if the next line should be grouped with current line

    Grouping logic:
    - Only group lines from the SAME character
    - Total length must be under MAX_CHARS_PER_SLIDE
    - Different characters always get separate slides
    """
    if not next_line:
        return False

    # Check if same character
    same_character = current_line['character'] == next_line['character']

    # NEVER group lines from different characters
    if not same_character:
        return False

    # Check total length
    next_length = len(next_line['translation'])
    total_length = current_length + next_length + 1  # +1 for space

    if total_length > MAX_CHARS_PER_SLIDE:
        return False

    # Group if same character and fits
    return True


def generate_ppt_from_json(json_file, output_file=None):
    """
    Generate PPT from translation JSON file

    Args:
        json_file: Path to translation JSON file
        output_file: Optional output PPT file path
    """
    # Load translation data
    with open(json_file, 'r', encoding='utf-8') as f:
        data = json.load(f)

    scene_id = data.get('sceneId', 'Unknown')
    scene_name = data.get('sceneName', '')
    translations = data['translations']

    print(f"Generating PPT for Scene {scene_id}: {scene_name}")
    print(f"Total lines: {len(translations)}")

    # Create presentation
    prs = create_presentation()

    # Add title slide
    title_slide_layout = prs.slide_layouts[6]
    title_slide = prs.slides.add_slide(title_slide_layout)
    background = title_slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

    title_box = title_slide.shapes.add_textbox(
        Inches(1), Inches(2),
        Inches(8), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = f"Scene {scene_id}\n{scene_name}"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = TEXT_COLOR
    title_para.alignment = PP_ALIGN.CENTER

    # Process translations with smart grouping
    i = 0
    slide_count = 0

    while i < len(translations):
        current = translations[i]
        grouped_text = current['translation']
        character = current['characterPinyin']
        current_length = len(grouped_text)

        # Try to group subsequent lines
        j = i + 1
        while j < len(translations):
            next_line = translations[j]
            if should_group_lines(current, next_line, current_length):
                grouped_text += "\n" + next_line['translation']
                current_length += len(next_line['translation']) + 1
                j += 1
            else:
                break

        # Split long text if needed
        text_parts = split_long_text(grouped_text)

        # Add slide(s) - one for each part
        for part in text_parts:
            add_subtitle_slide(prs, character, part, scene_id)
            slide_count += 1

        # Move to next ungrouped line
        i = j

    # Determine output filename
    if not output_file:
        base_name = os.path.splitext(os.path.basename(json_file))[0]
        output_file = f"{base_name}_subtitles.pptx"

    # Save presentation
    prs.save(output_file)
    print(f"\nGenerated {slide_count} subtitle slides")
    print(f"Saved to: {output_file}")

    return output_file


def main():
    """Main entry point"""
    if len(sys.argv) < 2:
        print("Usage: python3 generate_subtitles_ppt.py <translation_json_file>")
        print("Example: python3 generate_subtitles_ppt.py scene_1-4_translation.json")
        sys.exit(1)

    json_file = sys.argv[1]

    if not os.path.exists(json_file):
        print(f"Error: File not found: {json_file}")
        sys.exit(1)

    # Generate PPT
    try:
        output_file = generate_ppt_from_json(json_file)
        print(f"\n✓ Successfully generated: {output_file}")
    except Exception as e:
        print(f"\n✗ Error generating PPT: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == "__main__":
    main()
