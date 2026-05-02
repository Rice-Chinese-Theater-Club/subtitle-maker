#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Verify PPT content - check if all translations appear correctly in the generated PPTX.
"""

from pptx import Presentation
import json
import sys
import os


def verify(json_path, pptx_path):
    # Load translation data
    with open(json_path, 'r', encoding='utf-8') as f:
        data = json.load(f)

    translations = data['translations']

    # Load PPT
    prs = Presentation(pptx_path)

    print("=" * 80)
    print(f"PPT verification: {os.path.basename(pptx_path)}")
    print("=" * 80)
    print()

    # Skip title slide (index 0), check subtitle slides
    slides_to_check = min(11, len(prs.slides))

    for slide_idx in range(1, slides_to_check):
        slide = prs.slides[slide_idx]

        print(f"Slide {slide_idx}:")
        print("-" * 80)

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text = shape.text.strip()

                if len(text) < 10 and text.isupper():
                    print(f"  Character: {text}")
                elif text.startswith("Scene"):
                    print(f"  Scene: {text}")
                else:
                    print(f"  Content: {text}")

                    found = False
                    for idx, trans in enumerate(translations):
                        if text in trans['translation'] or trans['translation'] in text:
                            print(f"    -> Matches translation index {idx}: {trans['characterPinyin']} - {trans['original'][:30]}...")
                            found = True
                            break

                    if not found:
                        print(f"    -> No exact match found in translation JSON")

        print()

    print("=" * 80)
    print(f"Total slides: {len(prs.slides)} (including 1 title slide)")
    print("=" * 80)


if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python3 verify_ppt.py <translation_json> <output_pptx>")
        print("Example: python3 verify_ppt.py translations/scene_1-4.json outputs/scene_1-4_subtitles.pptx")
        sys.exit(1)
    verify(sys.argv[1], sys.argv[2])
